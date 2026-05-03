from __future__ import annotations

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader


class ExtractionError(Exception):
    pass


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            parts.append(t)
    return "\n\n".join(parts).strip()


def _iter_shape_text(shape) -> list[str]:
    out: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            out.extend(_iter_shape_text(child))
        return out
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                out.append(line)
    return out


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            parts.extend(_iter_shape_text(shape))
    return "\n".join(parts).strip()


def _find_soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    if sys.platform == "darwin":
        mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.isfile(mac):
            return mac
    return None


def _extract_ppt_via_libreoffice(path: Path) -> str:
    soffice = _find_soffice()
    if not soffice:
        raise ExtractionError(
            "旧版 .ppt 需安装 LibreOffice，或将文件在 PowerPoint / WPS 中另存为 .pptx 后上传。"
        )
    try:
        with tempfile.TemporaryDirectory() as td:
            td_path = Path(td)
            proc = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    str(td_path),
                    str(path.resolve()),
                ],
                check=False,
                timeout=120,
                capture_output=True,
                text=True,
            )
            if proc.returncode != 0:
                raise ExtractionError("LibreOffice 转换 .ppt 失败，请改为上传 .pptx。")
            out = td_path / (path.stem + ".pptx")
            if not out.exists():
                candidates = list(td_path.glob("*.pptx"))
                if not candidates:
                    raise ExtractionError("未能从 .ppt 得到 .pptx，请手动另存为 .pptx 后上传。")
                out = candidates[0]
            return _extract_pptx(out)
    except subprocess.TimeoutExpired as exc:
        raise ExtractionError("转换 .ppt 超时，请尝试更小的文件或改用 .pptx。") from exc


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        text = _extract_pdf(path)
    elif ext == ".pptx":
        text = _extract_pptx(path)
    elif ext == ".ppt":
        text = _extract_ppt_via_libreoffice(path)
    else:
        raise ExtractionError(f"不支持的扩展名：{ext}")

    if not text:
        raise ExtractionError("未能从文件中提取到文字，可能是扫描版 PDF 或空白演示文稿。")
    return text
